import logfire
from docx import Document
from pptx import Presentation

def parse_office(file_path: str):
    with logfire.span("Office Document Parsing", filename=file_path):
        try:
            extension = file_path.lower().rsplit(".", 1)[-1]
            if extension == "docx":
                document = Document(file_path)
                text_parts = [paragraph.text for paragraph in document.paragraphs]
                for table in document.tables:
                    text_parts.extend(cell.text for row in table.rows for cell in row.cells)
            elif extension == "pptx":
                presentation = Presentation(file_path)
                text_parts = [
                    shape.text
                    for slide in presentation.slides
                    for shape in slide.shapes
                    if hasattr(shape, "text")
                ]
            else:
                raise ValueError(f"Unsupported Office file type: .{extension}")

            full_text = "\n".join(text_parts)

            if not full_text.strip():
                logfire.warning(f"⚠️ Unstructured has return empty text for {file_path}")
            else:
                logfire.info(f"✔️ Successfully parse {len(full_text)} characters")

            return full_text
        except Exception as e:
            logfire.error(f"❌ Office Parsing Fail, {e}")
            raise e